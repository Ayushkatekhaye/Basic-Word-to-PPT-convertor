import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from textwrap import fill

def open_word_file():
    file_path = filedialog.askopenfilename(filetypes=[("Word files", "*.docx")])
    if file_path:
        ppt_file_path = filedialog.asksaveasfilename(defaultextension=".pptx")
        if ppt_file_path:
            word_to_ppt(file_path, ppt_file_path)

def word_to_ppt(word_file_path, ppt_file_path):
    ppt = Presentation()
    doc = Document(word_file_path)
    
    slide_layout = ppt.slide_layouts[5]
    background_image_path = "7243-01-low-poly-background-16x9-1.jpg"

    current_slide = None
    text_frame = None
    line_count = 0
    lines_per_slide = 15

    for paragraph in doc.paragraphs:
        if current_slide is None or line_count >= lines_per_slide:
            current_slide = ppt.slides.add_slide(slide_layout)
            current_slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
            text_frame = current_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4)).text_frame
            line_count = 0

        p = text_frame.add_paragraph()
        text = fill(paragraph.text, width=40)
        p.text = text
        line_count += len(text.split('\n'))

    if ppt.slides:
        ppt.save(ppt_file_path)
        messagebox.showinfo("Conversion Complete", "Word to PowerPoint conversion successful!")
    else:
        messagebox.showwarning("No Content", "No content to convert!")

window = tk.Tk()
window.title("Word to PowerPoint Converter")

window.geometry("800x500")
window.resizable(False, False)

label = tk.Label(window, text="Select a Word document to convert to PowerPoint:")
label.pack(pady=10)

browse_button = tk.Button(window, text="Browse", command=open_word_file)
browse_button.pack()

window.mainloop()